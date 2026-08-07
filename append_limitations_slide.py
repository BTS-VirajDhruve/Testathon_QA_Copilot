"""Append one Limitations slide to existing Agentic QA Copilot deck.

Does NOT recreate or modify existing slides — copies the PPTX, then adds slide 7.
"""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(r"D:\Projects\QA Copilot - Testathon")
SRC = ROOT / "Agentic_QA_Copilot_Demo.pptx"
DST = ROOT / "Agentic_QA_Copilot_Demo_With_Limitations.pptx"

# Theme matched to existing deck (slides 1-5 palette)
NAVY = RGBColor(0x0B, 0x12, 0x20)
CARD = RGBColor(0x16, 0x1F, 0x31)
CARD_ALT = RGBColor(0x15, 0x1C, 0x2E)
BORDER = RGBColor(0x2A, 0x3A, 0x55)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
ELECTRIC = RGBColor(0x38, 0xBD, 0xF8)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
GOLD = RGBColor(0xF5, 0xC5, 0x42)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CORAL = RGBColor(0xF8, 0x71, 0x71)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NOTES = (
    "This slide sets realistic expectations. The demo proves the core workflow: "
    "graph-based context, knowledge retrieval, AI generation, review, coverage gap "
    "detection, and evidence-backed outputs. The current limitations are mainly around "
    "production hardening: input completeness, modeled coverage boundaries, human "
    "approval, enterprise integrations, and runtime governance. These are exactly the "
    "areas planned for future roadmap expansion."
)

LIMITATIONS = [
    (
        "01",
        "Input quality matters",
        "QA output quality depends on the completeness of the system-flow graph and uploaded knowledge.",
        CORAL,
    ),
    (
        "02",
        "Coverage is modeled coverage",
        "Coverage is measured against modeled graph paths, requirements, risks, and evidence - not infinite real-world possibilities.",
        AMBER,
    ),
    (
        "03",
        "Human QA approval is still required",
        "Candidate bugs, generated tests, and automation recommendations should be reviewed by QA experts before release use.",
        GOLD,
    ),
    (
        "04",
        "Integration maturity",
        "Jira, Confluence, execution tools, and CI/CD need production-grade auth, sync, permissions, and monitoring.",
        ELECTRIC,
    ),
    (
        "05",
        "Runtime and cost governance",
        "LLM usage requires model routing, rate-limit handling, fallbacks, observability, and cost controls.",
        VIOLET,
    ),
]

ORBIT_NODES = [
    ("Input Quality", CORAL),
    ("Modeled Coverage", AMBER),
    ("Human Review", GOLD),
    ("Integrations", ELECTRIC),
    ("Runtime Governance", VIOLET),
]


def set_run(run, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(shape, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return tf


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def soft_border(shape, color=BORDER, width_pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def rounded_rect(slide, left, top, width, height, fill=CARD, border=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, fill)
    soft_border(shape, border, 1.25)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def oval(slide, left, top, width, height, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill_solid(shape, fill)
    if border is None:
        no_line(shape)
    else:
        soft_border(shape, border, 1.5)
    return shape


def accent_bar(slide, left, top, width, height, color=ELECTRIC):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(bar, color)
    no_line(bar)
    return bar


def sha256_slide_xml(prs_path: Path, slide_index: int) -> str:
    """Hash a slide part from the zip for change detection."""
    import zipfile

    # slide parts are ppt/slides/slideN.xml (1-based)
    name = f"ppt/slides/slide{slide_index + 1}.xml"
    with zipfile.ZipFile(prs_path, "r") as zf:
        data = zf.read(name)
    return hashlib.sha256(data).hexdigest()


def ensure_notes_text(prs: Presentation, slide, text: str) -> None:
    """Ensure notes body exists (this deck's blank notes slides lack placeholders)."""
    from copy import deepcopy
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from xml.sax.saxutils import escape

    notes = slide.notes_slide
    if notes.notes_text_frame is not None:
        notes.notes_text_frame.text = text
        return

    template = None
    for s in prs.slides:
        if s is slide:
            continue
        if s.has_notes_slide and s.notes_slide.notes_text_frame is not None:
            template = s.notes_slide
            break
    if template is None:
        raise RuntimeError("No template notes slide with body placeholder found")

    template_cSld = template._element.find(qn("p:cSld"))
    cloned = deepcopy(template_cSld)
    safe = escape(text)
    for sp in cloned.findall(".//" + qn("p:sp")):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        ph = nvPr.find(qn("p:ph"))
        if ph is not None and ph.get("type") == "body":
            txBody = sp.find(qn("p:txBody"))
            if txBody is not None:
                for child in list(txBody):
                    if child.tag == qn("a:p"):
                        txBody.remove(child)
                txBody.append(
                    parse_xml(
                        f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                        f"<a:r><a:t>{safe}</a:t></a:r></a:p>"
                    )
                )

    old_cSld = notes._element.find(qn("p:cSld"))
    notes._element.replace(old_cSld, cloned)
    # Text is already embedded in the cloned notes body XML.
    # Do not require notes_text_frame to resolve in-memory (shape cache quirks).
    body_texts = [
        t.text
        for t in notes._element.findall(".//" + qn("a:t"))
        if t.text
    ]
    if text not in body_texts and not any(text[:40] in (t or "") for t in body_texts):
        raise RuntimeError(f"Failed to embed notes text; found={body_texts!r}")


def add_limitations_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank)

    # Background + accent bars (match existing deck language)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(bg, NAVY)
    no_line(bg)
    accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), ELECTRIC)
    accent_bar(slide, 0, SLIDE_H - Inches(0.04), SLIDE_W, Inches(0.04), RGBColor(0x1E, 0x3A, 0x5F))

    # Title + subtitle (left)
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(7.2), Inches(0.55))
    add_text(title, "Limitations & Current Scope", size=28, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.82), Inches(7.4), Inches(0.4))
    add_text(
        sub,
        "What the current demo proves - and what remains for production hardening",
        size=13,
        color=SOFT,
    )
    accent_bar(slide, Inches(0.55), Inches(1.28), Inches(1.8), Inches(0.05), AMBER)

    # Main message
    msg = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(7.4), Inches(0.55))
    add_text(
        msg,
        "The demo validates the core agentic QA workflow, while some production capabilities require further hardening.",
        size=13,
        color=MUTED,
    )

    # Five limitation cards (left column)
    for i, (num, head, body, col) in enumerate(LIMITATIONS):
        top = Inches(2.05 + i * 0.78)
        card = rounded_rect(slide, Inches(0.55), top, Inches(7.5), Inches(0.7), fill=CARD)
        accent_bar(slide, Inches(0.55), top, Inches(0.07), Inches(0.7), col)

        nbox = slide.shapes.add_textbox(Inches(0.78), top + Inches(0.16), Inches(0.45), Inches(0.35))
        add_text(nbox, num, size=14, bold=True, color=col)

        hbox = slide.shapes.add_textbox(Inches(1.3), top + Inches(0.08), Inches(6.5), Inches(0.28))
        add_text(hbox, head, size=13, bold=True, color=WHITE)

        bbox = slide.shapes.add_textbox(Inches(1.3), top + Inches(0.35), Inches(6.5), Inches(0.32))
        add_text(bbox, body, size=11, color=MUTED)

    # ---- Right: scope-boundary / radar orbit visual ----
    cx = Inches(10.55)
    cy = Inches(3.55)

    # Outer radar rings
    for size, fill in [
        (Inches(4.4), RGBColor(0x10, 0x18, 0x28)),
        (Inches(3.35), RGBColor(0x12, 0x1C, 0x2E)),
        (Inches(2.25), RGBColor(0x14, 0x20, 0x34)),
    ]:
        oval(
            slide,
            cx - size / 2,
            cy - size / 2,
            size,
            size,
            fill=fill,
            border=RGBColor(0x2A, 0x3A, 0x55),
        )

    # Crosshair lines (subtle)
    hline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - Inches(2.1), cy - Pt(0.75), Inches(4.2), Pt(1.5)
    )
    fill_solid(hline, RGBColor(0x2A, 0x3A, 0x55))
    no_line(hline)
    vline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx - Pt(0.75), cy - Inches(2.1), Pt(1.5), Inches(4.2)
    )
    fill_solid(vline, RGBColor(0x2A, 0x3A, 0x55))
    no_line(vline)

    # Center core
    core_size = Inches(1.55)
    core = oval(
        slide,
        cx - core_size / 2,
        cy - core_size / 2,
        core_size,
        core_size,
        fill=RGBColor(0x12, 0x2A, 0x3A),
        border=TEAL,
    )
    core_label = slide.shapes.add_textbox(
        cx - Inches(0.85), cy - Inches(0.4), Inches(1.7), Inches(0.8)
    )
    tf = add_text(core_label, "Agentic QA", size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Copilot Core"
    set_run(run, size=12, bold=True, color=TEAL)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Demo"
    set_run(run2, size=11, bold=False, color=SOFT)

    # Orbiting limitation nodes (5 positions around circle)
    import math

    radius = Inches(1.85)
    node_w, node_h = Inches(1.55), Inches(0.55)
    for i, (label, col) in enumerate(ORBIT_NODES):
        # Start at top, go clockwise
        angle = -math.pi / 2 + i * (2 * math.pi / 5)
        nx = cx + int(radius * math.cos(angle)) - node_w // 2
        ny = cy + int(radius * math.sin(angle)) - node_h // 2

        # Connector dot on ring
        dot = oval(
            slide,
            cx + int(radius * math.cos(angle)) - Inches(0.08),
            cy + int(radius * math.sin(angle)) - Inches(0.08),
            Inches(0.16),
            Inches(0.16),
            fill=col,
        )

        node = rounded_rect(slide, nx, ny, node_w, node_h, fill=CARD_ALT, border=col)
        try:
            node.adjustments[0] = 0.2
        except Exception:
            pass
        ntxt = slide.shapes.add_textbox(nx, ny + Inches(0.12), node_w, Inches(0.35))
        add_text(ntxt, label, size=10, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bottom takeaway banner
    banner = rounded_rect(
        slide, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.85), fill=RGBColor(0x10, 0x18, 0x28)
    )
    soft_border(banner, TEAL, 1.25)

    take = slide.shapes.add_textbox(Inches(0.8), Inches(6.28), Inches(11.7), Inches(0.6))
    tf = add_text(
        take,
        "Current demo = context-aware QA intelligence.",
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Next step = continuous enterprise QA platform."
    set_run(run, size=14, bold=True, color=TEAL)

    # Slide number (match Future Scope style: "N  /  N")
    total = len(prs.slides)
    num_box = slide.shapes.add_textbox(Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.28))
    add_text(num_box, f"{total}  /  {total}", size=11, color=DIM, align=PP_ALIGN.RIGHT)

    # Speaker notes (do not touch existing slides' notes)
    ensure_notes_text(prs, slide, NOTES)


def extract_first_texts(prs: Presentation, max_per_slide: int = 3) -> list[list[str]]:
    out = []
    for s in prs.slides:
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = " ".join(sh.text_frame.text.split())
                if t:
                    texts.append(t[:80])
            if len(texts) >= max_per_slide:
                break
        out.append(texts)
    return out


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)

    # Exact file copy first so existing slide binaries start identical
    shutil.copy2(SRC, DST)

    # Hash original slide XMLs before mutation of the copy
    orig_hashes = [sha256_slide_xml(SRC, i) for i in range(len(Presentation(str(SRC)).slides))]
    orig_count = len(orig_hashes)
    orig_texts = extract_first_texts(Presentation(str(SRC)))

    prs = Presentation(str(DST))
    add_limitations_slide(prs)
    prs.save(str(DST))

    # Validation
    prs2 = Presentation(str(DST))
    new_count = len(prs2.slides)
    new_hashes = [sha256_slide_xml(DST, i) for i in range(orig_count)]
    unchanged = all(a == b for a, b in zip(orig_hashes, new_hashes))
    new_texts = extract_first_texts(prs2)
    text_unchanged = new_texts[:orig_count] == orig_texts

    # Confirm last slide
    last = prs2.slides[new_count - 1]
    last_title = ""
    for sh in last.shapes:
        if sh.has_text_frame and "Limitations" in sh.text_frame.text:
            last_title = sh.text_frame.text.strip()
            break
    notes_ok = False
    if last.has_notes_slide:
        notes_ok = "realistic expectations" in last.notes_slide.notes_text_frame.text

    print("INPUT:", SRC)
    print("OUTPUT:", DST)
    print("ORIG_SLIDES:", orig_count)
    print("NEW_SLIDES:", new_count)
    print("EXISTING_XML_UNCHANGED:", unchanged)
    print("EXISTING_TEXT_UNCHANGED:", text_unchanged)
    print("LAST_TITLE:", last_title)
    print("NOTES_ADDED:", notes_ok)
    print("OUTPUT_SIZE:", DST.stat().st_size)


if __name__ == "__main__":
    main()
