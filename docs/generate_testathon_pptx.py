"""Generate Agentic QA Copilot testathon PowerPoint with photo placeholders."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "Agentic_QA_Copilot_Testathon.pptx"

# Brand-ish palette (ink green / mist / brass — matches README design notes)
INK = RGBColor(0x1A, 0x2E, 0x28)
MIST = RGBColor(0xF4, 0xF7, 0xF5)
BRASS = RGBColor(0xB8, 0x9B, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x6B, 0x64)
ACCENT_BG = RGBColor(0xE8, 0xEF, 0xEB)
PLACEHOLDER_BG = RGBColor(0xDD, 0xE5, 0xE4)
PLACEHOLDER_BORDER = RGBColor(0x9A, 0xB0, 0xA6)


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, *, size=15, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)
    return box


def add_bg(slide, color=MIST):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back by z-order: python-pptx doesn't have send_to_back easily;
    # backgrounds are added first so content on top is fine.


def add_top_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.22), Inches(10), Inches(0.5), title, size=24, bold=True, color=WHITE)
    add_textbox(
        slide,
        Inches(10.2),
        Inches(0.28),
        Inches(2.8),
        Inches(0.4),
        "Agentic QA Copilot",
        size=12,
        color=BRASS,
        align=PP_ALIGN.RIGHT,
    )


def add_photo_placeholder(slide, left, top, width, height, label="Add screenshot / photo here"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_BG
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)
    # Dashed look isn't perfect in pptx; use text label instead
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.auto_size = None
    except Exception:
        pass
    # vertical center approx via empty paragraphs + main line
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "📷  PHOTO PLACEHOLDER"
    set_run(run, size=14, bold=True, color=MUTED)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = label
    set_run(run2, size=12, color=MUTED)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(6)
    run3 = p3.add_run()
    run3.text = "(replace this box with your image)"
    set_run(run3, size=10, color=PLACEHOLDER_BORDER)
    return shape


def feature_slide(prs, title, subtitle, bullets, photo_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    add_top_bar(slide, title)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.05), Inches(6.8), Inches(0.55), subtitle, size=14, color=MUTED)
    add_bullets(slide, Inches(0.5), Inches(1.65), Inches(6.8), Inches(5.2), bullets, size=16)
    add_photo_placeholder(
        slide,
        Inches(7.6),
        Inches(1.2),
        Inches(5.2),
        Inches(5.6),
        photo_label,
    )
    return slide


def two_photo_slide(prs, title, subtitle, bullets, photo_a, photo_b):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, title)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.45), subtitle, size=14, color=MUTED)
    add_bullets(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), bullets, size=15)
    add_photo_placeholder(slide, Inches(0.5), Inches(3.3), Inches(5.9), Inches(3.7), photo_a)
    add_photo_placeholder(slide, Inches(6.9), Inches(3.3), Inches(5.9), Inches(3.7), photo_b)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1), "Agentic QA Copilot", size=44, bold=True, color=WHITE)
    add_textbox(
        s,
        Inches(0.8),
        Inches(3.1),
        Inches(11.5),
        Inches(0.8),
        "Understand the software under test — then generate evidence-backed QA artifacts",
        size=20,
        color=BRASS,
    )
    add_textbox(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5), "Testathon Initiative Overview", size=16, color=MIST)
    add_textbox(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4), "Slide deck — add screenshots in the photo placeholders", size=12, color=MUTED)

    # 2. Problem
    feature_slide(
        prs,
        "The Problem",
        "Why generic AI test generation falls short",
        [
            "Chatbots invent tests that ignore real product flows and edge paths",
            "Coverage gaps (lockout, SSO timeout, payment failures) stay invisible",
            "Change impact is tribal knowledge — hard to brief regression quickly",
            "Little evidence trail for why a test exists or what it covers",
            "QA teams spend hours drafting suites that still miss high-risk branches",
        ],
        "Optional: photo of current pain / sticky notes / old process",
    )

    # 3. Solution overview
    feature_slide(
        prs,
        "What We Built",
        "Flow-aware, evidence-backed QA intelligence — not a chatbot",
        [
            "User-provided system flow graph as first-class product context",
            "Graph RAG + Vector RAG fused into one retrieval context",
            "Specialized agents: tests, exploratory, bugs, regression, impact, coverage",
            "Critic → coverage gaps → targeted regeneration loop",
            "Every artifact carries evidence + agent execution trace",
            "Works with OpenAI — or deterministic fallback for offline demos",
        ],
        "Hero UI screenshot — full app / landing",
    )

    # 4. Architecture
    feature_slide(
        prs,
        "How It Works",
        "End-to-end agentic pipeline",
        [
            "Query → load project flow → classify intent → plan retrieval",
            "Traverse graph paths + Vector RAG + existing tests/bugs",
            "Fuse context → specialist agents generate artifacts",
            "Critic reviews → prioritize high-risk gaps → targeted tests",
            "Return narrative + evidence + coverage before/after + trace",
            "Stack: Next.js frontend · FastAPI backend · OpenAI · Chroma · JSON graph",
        ],
        "Architecture diagram or pipeline screenshot",
    )

    # 5. System Flow Builder
    feature_slide(
        prs,
        "Feature: System Flow Builder",
        "Visual editor for the software under test",
        [
            "React Flow canvas — add/edit nodes and branches",
            "Import nested JSON or generate graph from natural language",
            "Undo/redo and export for sharing flows across the team",
            "Models auth methods, failure paths, alternate flows, components",
            "User-provided facts keep provenance — never silently overwritten by inference",
        ],
        "Screenshot: System Flow — Sign In branches",
    )

    # 6. Knowledge Base
    feature_slide(
        prs,
        "Feature: Knowledge Base",
        "Vector RAG over requirements and QA docs",
        [
            "Paste or upload requirements, notes, PDFs, DOCX",
            "Chunk → embed → store (Chroma, with JSON fallback)",
            "Optional entity extraction into the knowledge graph",
            "Retrieved at query time alongside graph paths",
            "Demo seeds authentication requirements; ecommerce pack includes billing/payment docs",
        ],
        "Screenshot: Knowledge Base panel with ingested docs",
    )

    # 7. Graph Explorer
    feature_slide(
        prs,
        "Feature: Graph Explorer",
        "Click any node for structured product insight",
        [
            "Dependencies, flows, linked tests and bugs",
            "Risk and coverage signals at the node level",
            "Makes Graph RAG tangible for demos and reviews",
            "Helps QA and PM walk the blast radius of a feature",
        ],
        "Screenshot: Graph Explorer node insight",
    )

    # 8. QA Copilot
    feature_slide(
        prs,
        "Feature: QA Copilot Console",
        "The main agentic workspace",
        [
            "Curated / suggested queries for common QA asks",
            "Optional changed-node for impact / regression focus",
            "Shows context used: flow paths, docs, tests, bugs",
            "Narrative answer + evidence cards + coverage regen metrics",
            "Labels LLM vs Deterministic fallback honestly",
        ],
        "Screenshot: Copilot before/during Run agentic analysis",
    )

    # 9. Test generation + evidence
    two_photo_slide(
        prs,
        "Feature: Evidence-Backed Test Generation",
        "Every test is grounded in graph path + sources",
        [
            "LLM-first structured generation with path-based deterministic fallback",
            "Cards explain why the test exists, path covered, and source evidence",
            "Security / negative / historical-bug focus in curated demo query",
        ],
        "Screenshot: generated test case card",
        "Screenshot: evidence / sources panel",
    )

    # 10. Critic + coverage loop
    feature_slide(
        prs,
        "Feature: Critic & Coverage Gap Loop",
        "Generate → review → find gaps → regenerate where risk is highest",
        [
            "Critic agent reviews initial artifacts",
            "Coverage analysis finds uncovered / weak branches",
            "Prioritizes critical/high gaps for targeted regen (bounded rounds)",
            "UI shows before → after coverage and remaining gaps",
            "Demo intentionally surfaces paths like Account Lockout / SSO Timeout",
        ],
        "Screenshot: coverage loop panel (before → after)",
    )

    # 11. Impact & Regression
    feature_slide(
        prs,
        "Feature: Impact & Regression Analysis",
        "Answer “what breaks if X changes?”",
        [
            "Select a changed node (e.g. Google OAuth) in Copilot",
            "Impact agent walks the subgraph and dependent components",
            "Regression agent recommends a focused suite",
            "Turns tribal knowledge into a repeatable briefing",
        ],
        "Screenshot: impact / regression results",
    )

    # 12. Artifacts catalog
    feature_slide(
        prs,
        "Feature: Artifact Catalogs",
        "Tests · Exploratory · Bugs · Regression — in one place",
        [
            "Seeded + generated test cases browsable in the UI",
            "Exploratory missions from Copilot runs",
            "Historical and generated bug reports",
            "Regression recommendations after change analysis",
            "Dashboard cards: risk, coverage %, counts, confidence",
        ],
        "Screenshot: Test Cases / Bugs / Dashboard cards",
    )

    # 13. Agent Trace
    feature_slide(
        prs,
        "Feature: Agent Trace & Transparency",
        "See what the system actually did",
        [
            "Step-by-step orchestrator execution",
            "Skipped steps remain visible (honest pipeline view)",
            "Model routing diagnostics: which model, escalation, fallback",
            "Builds trust for enterprise / audit-minded stakeholders",
        ],
        "Screenshot: Agent Trace panel",
    )

    # 14. Demo projects
    two_photo_slide(
        prs,
        "Demo Projects",
        "Two ways to show the system live",
        [
            "One-click: Enterprise Authentication Portal (Sign In) — email, Google OAuth, Microsoft SSO, self-registration",
            "Manual pack: ShopEase ecommerce — catalog, cart, discounts, GST/billing, card/UPI/COD payments",
        ],
        "Screenshot: Load Demo Project / Sign In flow",
        "Screenshot: ShopEase flow or sample_data pack",
    )

    # 15. Tools used
    feature_slide(
        prs,
        "Tools Used",
        "What helped us get here",
        [
            "Cursor — AI-assisted architecture and implementation",
            "OpenAI API — chat + embeddings (task-routed models)",
            "Next.js / React / Tailwind / React Flow — UI",
            "FastAPI / uv / Pydantic — agentic backend",
            "ChromaDB + durable JSON graph (+ optional Neo4j sync)",
            "Docker Compose · pytest for packaging and tests",
        ],
        "Optional: team / IDE / architecture whiteboard photo",
    )

    # 16. Pending work
    feature_slide(
        prs,
        "Pending Work",
        "Path from prototype → team-ready product",
        [
            "Live external research (CVE/OWASP-style) — flagged but not executed",
            "Neo4j as active query path; production graph/vector backends",
            "Auth, multi-user tenancy, shared workspaces",
            "One-click ecommerce seed + ALM export (Jira / ADO / TestRail)",
            "Enable/tune premium reviewer; production hardening & observability",
            "Richer regen strategies + human-in-the-loop approval",
        ],
        "Optional: roadmap whiteboard photo",
    )

    # 17. Value / ROI
    feature_slide(
        prs,
        "Value Add & ROI",
        "Our assessment for the team",
        [
            "Faster first drafts of security/negative/regression suites",
            "Higher signal: coverage moves on highest-risk uncovered paths",
            "Change-impact briefings reduce regression guesswork",
            "Evidence + traces support auditability and review",
            "Low pilot risk: deterministic fallback + flow/docs import",
            "Largest ROI unlock after ALM sync + shared projects (pending)",
        ],
        "Optional: metrics / before-after collage",
    )

    # 18. Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_textbox(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1), "Not a chatbot.", size=40, bold=True, color=WHITE)
    add_textbox(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11.5),
        Inches(1),
        "Flow-aware · Evidence-backed · Self-improving coverage",
        size=22,
        color=BRASS,
    )
    add_textbox(
        s,
        Inches(0.8),
        Inches(5.0),
        Inches(11.5),
        Inches(0.6),
        "Agentic QA Copilot — Testathon",
        size=16,
        color=MIST,
    )
    add_photo_placeholder(s, Inches(8.5), Inches(4.8), Inches(4.2), Inches(2.2), "Team / demo photo (optional)")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
