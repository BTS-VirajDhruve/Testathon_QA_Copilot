"""Create Agentic QA Copilot 5-slide professional demo deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Theme colors
NAVY = RGBColor(0x0B, 0x12, 0x20)
CHARCOAL = RGBColor(0x11, 0x18, 0x27)
CARD = RGBColor(0x16, 0x1F, 0x31)
CARD_ALT = RGBColor(0x1A, 0x24, 0x38)
BORDER = RGBColor(0x2A, 0x3A, 0x55)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
ELECTRIC = RGBColor(0x38, 0xBD, 0xF8)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
GOLD = RGBColor(0xF5, 0xC5, 0x42)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CORAL = RGBColor(0xF8, 0x71, 0x71)

ROOT = Path(__file__).resolve().parent
ASSETS = Path(r"C:\Users\Viraj.Dhruve\.cursor\projects\d-Projects-QA-Copilot-Testathon\assets")
OUT_PPTX = ROOT / "Agentic_QA_Copilot_Demo.pptx"
OUT_PDF = ROOT / "Agentic_QA_Copilot_Demo.pdf"


def set_run(run, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(shape, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=14, bold=False, color=SOFT, align=PP_ALIGN.LEFT, space_before=6, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def soft_border(shape, color=BORDER, width_pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def rounded_rect(slide, left, top, width, height, fill=CARD, border=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(shape, fill)
    soft_border(shape, border, 1.25)
    # Slightly tighter corner
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def accent_bar(slide, left, top, width, height, color=ELECTRIC):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(bar, color)
    no_line(bar)
    return bar


def circle(slide, left, top, size, fill=ELECTRIC):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    fill_solid(shape, fill)
    no_line(shape)
    return shape


def set_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.name = "Calibri"


def dark_bg(slide):
    """Solid dark navy background + subtle top accent gradient strip."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(bg, NAVY)
    no_line(bg)
    # Top glow strip
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    fill_solid(top, ELECTRIC)
    no_line(top)
    # Bottom thin accent
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.04), SLIDE_W, Inches(0.04))
    fill_solid(bottom, RGBColor(0x1E, 0x3A, 0x5F))
    no_line(bottom)
    return bg


def add_image_dimmed(slide, path, left, top, width, height):
    """Place hero image with a dark overlay for text readability."""
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(overlay, RGBColor(0x08, 0x0F, 0x1C))
    # Semi-transparent via alpha on solid fill
    spPr = overlay.fill._xPr
    solidFill = spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    if solidFill is not None:
        srgb = solidFill.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
        )
        if srgb is not None:
            alpha = parse_xml(
                '<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="45000"/>'
            )
            srgb.append(alpha)
    no_line(overlay)
    return pic


def pill(slide, left, top, width, height, text, fill=RGBColor(0x1E, 0x3A, 0x5F), color=ELECTRIC):
    shape = rounded_rect(slide, left, top, width, height, fill=fill, border=fill)
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=11, bold=True, color=color)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    # Vertically center-ish
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)
    return shape


# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------

def slide_1_title(prs, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    dark_bg(slide)

    # Full-bleed hero with overlay
    add_image_dimmed(slide, img, 0, 0, SLIDE_W, SLIDE_H)

    # Re-draw top accent over image
    accent_bar(slide, 0, 0, SLIDE_W, Inches(0.06), ELECTRIC)

    # Brand label
    label = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(8), Inches(0.4))
    add_text(label, "LIVE DEMO", size=13, bold=True, color=TEAL)

    # Title
    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.1))
    add_text(title, "Agentic QA Copilot", size=48, bold=True, color=WHITE)

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.15), Inches(11), Inches(0.55))
    add_text(
        sub,
        "Graph RAG + Vector RAG + AI Agents for Evidence-Backed QA",
        size=20,
        bold=False,
        color=ELECTRIC,
    )

    # Accent underline
    accent_bar(slide, Inches(0.9), Inches(3.85), Inches(2.2), Inches(0.06), TEAL)

    # Key message
    msg = slide.shapes.add_textbox(Inches(0.9), Inches(4.15), Inches(10.5), Inches(1.2))
    tf = add_text(
        msg,
        "A QA assistant that understands the software flow, reads project knowledge,",
        size=16,
        color=SOFT,
    )
    add_para(
        tf,
        "generates test coverage, reviews gaps, and improves the suite automatically.",
        size=16,
        color=SOFT,
        space_before=4,
    )

    # Bottom tags
    for i, (txt, col) in enumerate(
        [
            ("System Flow Graph", TEAL),
            ("Knowledge Fusion", ELECTRIC),
            ("Agentic Coverage", VIOLET),
        ]
    ):
        pill(
            slide,
            Inches(0.9 + i * 2.6),
            Inches(6.35),
            Inches(2.4),
            Inches(0.42),
            txt,
            fill=RGBColor(0x12, 0x1A, 0x2B),
            color=col,
        )

    set_notes(
        slide,
        "Welcome the audience. Introduce Agentic QA Copilot as a live demo of an "
        "evidence-backed QA assistant. Emphasize that it understands product flow, "
        "fuses project knowledge, and uses agents to generate, review, and improve "
        "test coverage - not just prompt-based generation.",
    )
    return slide


def slide_2_problem(prs, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    # Right-side hero visual
    add_image_dimmed(slide, img, Inches(8.2), Inches(1.2), Inches(4.7), Inches(5.5))
    # Frame
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.2), Inches(4.7), Inches(5.5)
    )
    frame.fill.background()
    soft_border(frame, RGBColor(0x33, 0x45, 0x68), 1.5)
    try:
        frame.adjustments[0] = 0.04
    except Exception:
        pass

    # Title
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    add_text(title, "Why Traditional AI Test Generation Falls Short", size=28, bold=True)

    accent_bar(slide, Inches(0.7), Inches(1.15), Inches(1.8), Inches(0.05), AMBER)

    pains = [
        (
            "01",
            "Generic test cases",
            "AI tools often generate broad tests without understanding the actual product flow.",
            CORAL,
        ),
        (
            "02",
            "Missing business context",
            "Requirements, bugs, risks, and existing tests are scattered across documents and tools.",
            AMBER,
        ),
        (
            "03",
            "Weak coverage visibility",
            "QA teams struggle to see what is covered, what is missing, and what should be tested next.",
            GOLD,
        ),
    ]

    for i, (num, head, body, accent) in enumerate(pains):
        top = Inches(1.55 + i * 1.7)
        card = rounded_rect(slide, Inches(0.7), top, Inches(7.1), Inches(1.5), fill=CARD)
        accent_bar(slide, Inches(0.7), top, Inches(0.08), Inches(1.5), accent)

        num_box = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.28), Inches(0.8), Inches(0.45))
        add_text(num_box, num, size=22, bold=True, color=accent)

        h = slide.shapes.add_textbox(Inches(1.8), top + Inches(0.28), Inches(5.6), Inches(0.4))
        add_text(h, head, size=18, bold=True, color=WHITE)

        b = slide.shapes.add_textbox(Inches(1.8), top + Inches(0.72), Inches(5.6), Inches(0.6))
        add_text(b, body, size=13, color=MUTED)

    set_notes(
        slide,
        "Explain that normal prompt-based test generation is not enough because it does "
        "not understand product structure. Call out the three gaps: generic tests, "
        "scattered context, and weak coverage visibility. Set up the need for graph + "
        "knowledge-aware agents.",
    )
    return slide


def slide_3_architecture(prs, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.6))
    add_text(title, "How Our QA Copilot Understands the Product", size=28, bold=True)
    accent_bar(slide, Inches(0.7), Inches(1.0), Inches(1.8), Inches(0.05), TEAL)

    # Subtle background visual on right
    add_image_dimmed(slide, img, Inches(8.4), Inches(1.25), Inches(4.5), Inches(3.2))

    # Pipeline stages as horizontal cards
    sources = [
        ("System Flow Graph", TEAL),
        ("Knowledge Base", ELECTRIC),
        ("Jira / Confluence", VIOLET),
    ]
    for i, (label, col) in enumerate(sources):
        left = Inches(0.7 + i * 2.45)
        card = rounded_rect(slide, left, Inches(1.35), Inches(2.3), Inches(0.7), fill=CARD_ALT)
        add_text(card, label, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        card.text_frame.margin_top = Pt(14)

    # Arrow down 1
    arrow1 = slide.shapes.add_textbox(Inches(3.6), Inches(2.1), Inches(2), Inches(0.35))
    add_text(arrow1, "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Graph RAG + Vector RAG
    rag = rounded_rect(slide, Inches(1.6), Inches(2.45), Inches(5.8), Inches(0.7), fill=RGBColor(0x12, 0x2A, 0x3A))
    soft_border(rag, TEAL, 1.5)
    add_text(rag, "Graph RAG  +  Vector RAG", size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    rag.text_frame.margin_top = Pt(14)

    arrow2 = slide.shapes.add_textbox(Inches(3.6), Inches(3.2), Inches(2), Inches(0.35))
    add_text(arrow2, "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Context Fusion
    fusion = rounded_rect(slide, Inches(2.2), Inches(3.55), Inches(4.6), Inches(0.7), fill=RGBColor(0x1A, 0x18, 0x35))
    soft_border(fusion, VIOLET, 1.5)
    add_text(fusion, "Context Fusion", size=16, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    fusion.text_frame.margin_top = Pt(14)

    arrow3 = slide.shapes.add_textbox(Inches(3.6), Inches(4.3), Inches(2), Inches(0.35))
    add_text(arrow3, "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # QA Agents
    agents = rounded_rect(slide, Inches(2.6), Inches(4.65), Inches(3.8), Inches(0.7), fill=RGBColor(0x1A, 0x2A, 0x1F))
    soft_border(agents, GOLD, 1.5)
    add_text(agents, "QA Agents", size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    agents.text_frame.margin_top = Pt(14)

    # Explanation bullets on bottom / right
    explanations = [
        ("Graph RAG", "Workflows, branches, dependencies, failure paths", TEAL),
        ("Vector RAG", "Requirements, bugs, risks, and documents", ELECTRIC),
        ("Context Fusion", "Combines both into high-quality prompt context", VIOLET),
        ("QA Agents", "Generate, review, and refine QA outputs", GOLD),
    ]

    for i, (head, body, col) in enumerate(explanations):
        left = Inches(0.7 + (i % 2) * 6.2)
        top = Inches(5.6 + (i // 2) * 0.75)
        # small accent dot
        circle(slide, left, top + Inches(0.12), Inches(0.18), fill=col)
        hbox = slide.shapes.add_textbox(left + Inches(0.35), top, Inches(5.5), Inches(0.3))
        tf = add_text(hbox, head, size=13, bold=True, color=WHITE)
        add_para(tf, body, size=11, color=MUTED, space_before=0)

    set_notes(
        slide,
        "Walk through the architecture simply: the system flow graph plus knowledge "
        "sources feed Graph RAG and Vector RAG. Context Fusion merges both, then QA "
        "agents use that rich context to generate and refine outputs. Graph RAG "
        "understands structure; Vector RAG retrieves documents; together they enable "
        "evidence-backed QA.",
    )
    return slide


def slide_4_workflow(prs, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(12), Inches(0.55))
    add_text(title, "From Input to Complete QA Output", size=28, bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.9), Inches(1.8), Inches(0.05), ELECTRIC)

    # Background accent image (subtle, bottom-right)
    add_image_dimmed(slide, img, Inches(9.0), Inches(4.6), Inches(4.0), Inches(2.6))

    steps = [
        "Build system flow graph",
        "Add knowledge / import Jira & Confluence",
        "Run agentic analysis",
        "Generate test cases",
        "Review validity",
        "Classify automation feasibility",
        "Detect coverage gaps",
        "Generate targeted tests",
        "Export BDD / Gherkin",
    ]

    # 3x3 step grid
    for i, step in enumerate(steps):
        row, col = divmod(i, 3)
        left = Inches(0.7 + col * 2.85)
        top = Inches(1.2 + row * 0.95)
        card = rounded_rect(slide, left, top, Inches(2.7), Inches(0.82), fill=CARD)

        num = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.18), Inches(0.4), Inches(0.4))
        add_text(num, f"{i+1}", size=16, bold=True, color=ELECTRIC)

        txt = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.18), Inches(2.05), Inches(0.55))
        add_text(txt, step, size=12, bold=False, color=SOFT)

    # Loop banner
    loop = rounded_rect(
        slide, Inches(0.7), Inches(4.2), Inches(8.3), Inches(0.55), fill=RGBColor(0x12, 0x1A, 0x2B)
    )
    soft_border(loop, ELECTRIC, 1.25)
    add_text(
        loop,
        "Agent Loop   ·   Generate  →  Review  →  Fix  →  Cover  →  Export",
        size=13,
        bold=True,
        color=ELECTRIC,
        align=PP_ALIGN.CENTER,
    )
    loop.text_frame.margin_top = Pt(10)

    # Outputs section
    out_label = slide.shapes.add_textbox(Inches(0.7), Inches(4.95), Inches(4), Inches(0.35))
    add_text(out_label, "OUTPUTS", size=11, bold=True, color=GOLD)

    outputs = [
        "Test Cases",
        "BDD Feature Files",
        "Automation Review",
        "Exploratory Missions",
        "Bug Reports",
        "Regression Recs",
        "Coverage Gaps",
        "Sources & Evidence",
        "Agent Trace",
    ]

    for i, out in enumerate(outputs):
        left = Inches(0.7 + (i % 5) * 1.7)
        top = Inches(5.4 + (i // 5) * 0.55)
        pill(
            slide,
            left,
            top,
            Inches(1.55),
            Inches(0.4),
            out,
            fill=RGBColor(0x15, 0x1E, 0x2E),
            color=SOFT,
        )

    set_notes(
        slide,
        "Walk through the nine steps as the demo path. Emphasize the closed agent "
        "loop: generate, review, fix, cover, export. Highlight the breadth of outputs: "
        "tests, BDD, automation feasibility, coverage gaps, and evidence traces. "
        "Keep it light; the live demo will show these in action.",
    )
    return slide


def slide_5_impact(prs, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    # Soft hero on far right
    add_image_dimmed(slide, img, Inches(9.3), Inches(1.1), Inches(3.7), Inches(4.2))

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.55))
    add_text(title, "What Makes This Demo Powerful", size=28, bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.95), Inches(1.8), Inches(0.05), GOLD)

    cards = [
        ("01", "Context-aware QA", "Uses actual system flow and project knowledge.", TEAL),
        ("02", "Evidence-backed results", "Every output links to graph paths, requirements, bugs, or documents.", ELECTRIC),
        ("03", "Closed-loop coverage", "Reviewer finds weak or missing tests and sends them back for refinement.", VIOLET),
        ("04", "Automation-ready planning", "Classifies tests as automatable, manual, hybrid, or not ready.", GOLD),
    ]

    for i, (num, head, body, col) in enumerate(cards):
        row, col_i = divmod(i, 2)
        left = Inches(0.7 + col_i * 4.2)
        top = Inches(1.3 + row * 1.7)
        card = rounded_rect(slide, left, top, Inches(4.0), Inches(1.5), fill=CARD)
        accent_bar(slide, left, top, Inches(0.08), Inches(1.5), col)

        n = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), Inches(0.7), Inches(0.35))
        add_text(n, num, size=16, bold=True, color=col)

        h = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55), Inches(3.4), Inches(0.35))
        add_text(h, head, size=15, bold=True, color=WHITE)

        b = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.95), Inches(3.4), Inches(0.4))
        add_text(b, body, size=12, color=MUTED)

    # Live demo focus banner
    banner = rounded_rect(
        slide, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.7), fill=RGBColor(0x10, 0x18, 0x28)
    )
    soft_border(banner, TEAL, 1.5)

    focus_label = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11), Inches(0.35))
    add_text(focus_label, "LIVE DEMO FOCUS", size=12, bold=True, color=TEAL)

    focus_flow = slide.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(11.2), Inches(0.7))
    add_text(
        focus_flow,
        "System Flow  →  Knowledge Import  →  Agentic Analysis  →  Results  →  BDD Export",
        size=18,
        bold=True,
        color=WHITE,
    )

    set_notes(
        slide,
        "Close with the four differentiators: context-aware, evidence-backed, closed-loop "
        "coverage, and automation-ready planning. Then transition into the live demo using "
        "the focus path: System Flow → Knowledge Import → Agentic Analysis → Results → BDD Export.",
    )
    return slide


def export_pdf_via_powerpoint(pptx_path: Path, pdf_path: Path) -> bool:
    """Try Windows PowerPoint COM export to PDF."""
    try:
        import win32com.client  # type: ignore
    except ImportError:
        try:
            import comtypes.client as cc  # type: ignore

            powerpoint = cc.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
            # 32 = ppSaveAsPDF
            deck.SaveAs(str(pdf_path), 32)
            deck.Close()
            powerpoint.Quit()
            return pdf_path.exists()
        except Exception as e:
            print(f"PDF export via comtypes failed: {e}")
            return False

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.SaveAs(str(pdf_path), 32)
        deck.Close()
        powerpoint.Quit()
        return pdf_path.exists()
    except Exception as e:
        print(f"PDF export via win32com failed: {e}")
        return False


def main():
    imgs = {
        1: ASSETS / "slide1_hero.png",
        2: ASSETS / "slide2_hero.png",
        3: ASSETS / "slide3_hero.png",
        4: ASSETS / "slide4_hero.png",
        5: ASSETS / "slide5_hero.png",
    }
    for p in imgs.values():
        if not p.exists():
            raise FileNotFoundError(f"Missing hero image: {p}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs, imgs[1])
    slide_2_problem(prs, imgs[2])
    slide_3_architecture(prs, imgs[3])
    slide_4_workflow(prs, imgs[4])
    slide_5_impact(prs, imgs[5])

    prs.save(str(OUT_PPTX))
    print(f"Saved PPTX: {OUT_PPTX}")
    print(f"Slide count: {len(prs.slides)}")

    pdf_ok = export_pdf_via_powerpoint(OUT_PPTX, OUT_PDF)
    if pdf_ok:
        print(f"Saved PDF: {OUT_PDF}")
    else:
        print("PDF export skipped or failed (PowerPoint COM unavailable).")


if __name__ == "__main__":
    main()
